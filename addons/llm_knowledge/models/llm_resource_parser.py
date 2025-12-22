import base64
import json
import logging

from markdownify import markdownify as md

try:
    import pymupdf
except ImportError:
    pymupdf = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

from odoo import _, api, fields, models
from odoo.exceptions import UserError

_logger = logging.getLogger(__name__)


class LLMResourceParser(models.Model):
    _inherit = "llm.resource"

    parser = fields.Selection(
        selection="_get_available_parsers",
        string="Parser",
        default="default",
        required=True,
        help="Method used to parse resource content",
        tracking=True,
    )

    @api.model
    def _get_available_parsers(self):
        """Get all available parser methods"""
        return [
            ("default", "Default Parser"),
            ("json", "JSON Parser"),
            ("ocr", "OCR Parser (for images)"),
        ]

    def parse(self):
        """Parse the retrieved content to markdown"""
        # Lock resources and process only the successfully locked ones
        resources = self._lock(state_filter="retrieved")
        if not resources:
            return False

        for resource in resources:
            try:
                # Get the related record
                record = self.env[resource.res_model].browse(resource.res_id)
                if not record.exists():
                    raise UserError(_("Referenced record not found"))

                # If the record has a specific rag_parse method, call it
                if hasattr(record, "llm_get_fields"):
                    fields = record.llm_get_fields(record)
                else:
                    # Call get_fields on the individual resource to ensure singleton
                    fields = resource.get_fields(record)

                for field in fields:
                    # TODO: Should it be self._parse_field?
                    success = resource._parse_field(record, field)

                if success:
                    resource.write({"state": "parsed"})
                    self.env.cr.commit()
                    resource._post_styled_message(
                        "Resource successfully parsed", "success"
                    )
                else:
                    resource._post_styled_message(
                        "Parsing completed but did not return success", "warning"
                    )

            except Exception as e:
                _logger.error(
                    "Error parsing resource %s: %s",
                    resource.id,
                    str(e),
                    exc_info=True,
                )
                resource._post_styled_message(
                    f"Error parsing resource: {str(e)}", "error"
                )
                if resource.collection_ids:
                    resource.collection_ids._post_styled_message(
                        f"Error parsing resource: {str(e)}", "error"
                    )
            finally:
                resource._unlock()
        resources._unlock()

    def _get_parser(self, record, field_name, mimetype):
        # Check for explicit parser selection
        if self.parser == "ocr":
            return self._parse_image_ocr
        elif self.parser != "default":
            return getattr(self, f"parse_{self.parser}")
        record_name = (
            record.display_name
            if hasattr(record, "display_name")
            else f"{record._name} #{record.id}"
        )

        is_markdown = ".md" in record_name.lower()
        if mimetype == "application/pdf":
            return self._parse_pdf
        # special case, as odoo detects markdowns as application/octet-stream
        elif mimetype == "application/octet-stream" and is_markdown:
            return self._parse_text
        elif "html" in mimetype:
            return self._parse_html
        elif mimetype.startswith("text/"):
            return self._parse_text
        elif mimetype.startswith("image/"):
            # For images, use OCR to extract text (if available)
            return self._parse_image_ocr
        elif mimetype == "application/json":
            return self.parse_json
        elif mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self._parse_docx
        elif mimetype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self._parse_xlsx
        elif mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._parse_pptx
        elif mimetype == "text/csv" or (mimetype == "application/octet-stream" and record_name.lower().endswith(".csv")):
            return self._parse_csv
        else:
            return self._parse_default

    def _parse_field(self, record, field):
        self.ensure_one()
        parser_method = self._get_parser(record, field["field_name"], field["mimetype"])
        return parser_method(record, field)

    def get_fields(self, record):
        """
        Default parser implementation - generates a generic markdown representation
        based on commonly available fields

        :returns
        [{"field_name": field_name, "mimetype": mimetype, "rawcontent": rawcontent}]
        """
        self.ensure_one()

        results = []

        # Start with the record name/display_name if available
        record_name_field = (
            "display_name" if hasattr(record, "display_name") else "name"
        )
        record_name = (
            record[record_name_field]
            if hasattr(record, record_name_field)
            else f"{record._name} #{record.id}"
        )
        if record_name:
            results.append(
                {
                    "field_name": record_name_field,
                    "mimetype": "text/plain",
                    "rawcontent": record_name,
                }
            )

        # Try to include description or common text fields
        common_text_fields = [
            "description",
            "note",
            "comment",
            "message",
            "content",
            "body",
            "text",
        ]
        for field_name in common_text_fields:
            if hasattr(record, field_name) and record[field_name]:
                # Use text/plain for now, could be refined based on field type
                results.append(
                    {
                        "field_name": field_name,
                        "mimetype": "text/plain",
                        "rawcontent": record[field_name],
                    }
                )

        return results

    def parse_json(self, record, field):
        """
        JSON parser implementation - converts record data to JSON and then to markdown
        """
        self.ensure_one()

        # Get record name or default to model name and ID
        record_name = (
            record.display_name
            if hasattr(record, "display_name")
            else f"{record._name} #{record.id}"
        )

        # Create a dictionary with record data
        record_data = {}
        for field_name, field in record._fields.items():
            try:
                # Skip binary fields and internal fields
                if field.type == "binary" or field_name.startswith("_"):
                    continue

                # Handle many2one fields
                if field.type == "many2one" and record[field_name]:
                    record_data[field_name] = {
                        "id": record[field_name].id,
                        "name": record[field_name].display_name,
                    }
                # Handle many2many and one2many fields
                elif field.type in ["many2many", "one2many"]:
                    record_data[field_name] = [
                        {"id": r.id, "name": r.display_name} for r in record[field_name]
                    ]
                # Handle other fields
                else:
                    record_data[field_name] = record[field_name]
            except Exception as e:
                _logger.error(f"Skipping field {field_name}: {str(e)}")
                self._post_styled_message(
                    f"Skipping field {field_name}: {str(e)}", "warning"
                )
                continue
        # Format as markdown
        content = [f"# {record_name}"]
        content.append("\n## JSON Data\n")
        content.append("```json")
        content.append(json.dumps(record_data, indent=2, default=str))
        content.append("```")

        # Update resource content
        self.content = "\n".join(content)

        return True

    def _parse_pdf(self, record, field):
        """Parse PDF file and extract text and images"""
        # Decode attachment data

        if field["mimetype"] != "application/pdf":
            return False

        # Open PDF using PyMuPDF
        text_content = []
        image_count = 0
        page_count = 0
        # no need to decode as passing raw data should work here
        pdf_data = field["rawcontent"]

        # Create a BytesIO object from the PDF data
        with pymupdf.open(stream=pdf_data, filetype="pdf") as doc:
            # Store page count before document is closed
            page_count = doc.page_count

            # Process each page
            for page_num in range(page_count):
                page = doc[page_num]

                # Extract text
                text = page.get_text()
                text_content.append(f"## Page {page_num + 1}\n\n{text}")

                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    try:
                        base_image = doc.extract_image(xref)
                        if base_image:
                            # Store image as attachment
                            image_data = base_image["image"]
                            image_ext = base_image["ext"]
                            image_name = f"image_{page_num}_{img_index}.{image_ext}"

                            # Create attachment for the image
                            img_attachment = record.env["ir.attachment"].create(
                                {
                                    "name": image_name,
                                    "datas": base64.b64encode(image_data),
                                    "res_model": "llm.resource",
                                    "res_id": self.id,
                                    "mimetype": f"image/{image_ext}",
                                }
                            )

                            # Add image reference to markdown content
                            if img_attachment:
                                image_url = f"/web/image/{img_attachment.id}"
                                text_content.append(f"\n![{image_name}]({image_url})\n")
                                image_count += 1
                    except Exception as e:
                        self._post_styled_message(
                            f"Error extracting image: {str(e)}", "warning"
                        )

        # Join all content
        final_content = "\n\n".join(text_content)

        # Update resource with extracted content
        self.content = final_content

        return True

    def _parse_text(self, _, field):
        self.content = field["rawcontent"]
        return True

    def _parse_html(self, _, field):
        self.content = md(field["rawcontent"])
        return True

    def _parse_image(self, record, _):
        image_url = f"/web/image/{record.id}"
        self.content = f"![{record.name}]({image_url})"
        return True

    def _parse_docx(self, record, field):
        """Parse DOCX file and extract text content"""
        if DocxDocument is None:
            _logger.warning("python-docx not installed, falling back to default parser")
            return self._parse_default(record, field)

        import io
        docx_data = field["rawcontent"]

        try:
            # Create BytesIO from raw data
            docx_stream = io.BytesIO(docx_data)
            doc = DocxDocument(docx_stream)

            # Extract text from paragraphs
            text_content = []
            text_content.append(f"# {record.name}\n")

            for para in doc.paragraphs:
                if para.text.strip():
                    # Check if it's a heading
                    if para.style and para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        try:
                            level_num = int(level)
                            text_content.append(f"{'#' * (level_num + 1)} {para.text}")
                        except ValueError:
                            text_content.append(f"## {para.text}")
                    else:
                        text_content.append(para.text)

            # Extract text from tables
            for table in doc.tables:
                table_md = []
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    table_md.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        table_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                text_content.append("\n" + "\n".join(table_md) + "\n")

            self.content = "\n\n".join(text_content)
            return True

        except Exception as e:
            _logger.error(f"Error parsing DOCX: {str(e)}")
            self._post_styled_message(f"Error parsing DOCX: {str(e)}", "error")
            return self._parse_default(record, field)

    def _parse_xlsx(self, record, field):
        """Parse Excel (.xlsx) file and extract content as markdown tables"""
        if openpyxl is None:
            _logger.warning("openpyxl not installed, falling back to default parser")
            return self._parse_default(record, field)

        import io
        xlsx_data = field["rawcontent"]

        try:
            xlsx_stream = io.BytesIO(xlsx_data)
            workbook = openpyxl.load_workbook(xlsx_stream, read_only=True, data_only=True)

            text_content = [f"# {record.name}\n"]

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"\n## Sheet: {sheet_name}\n")

                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    text_content.append("*Empty sheet*\n")
                    continue

                # Find max columns with data
                max_cols = max(len([c for c in row if c is not None]) for row in rows) if rows else 0
                if max_cols == 0:
                    text_content.append("*Empty sheet*\n")
                    continue

                # Build markdown table
                table_md = []
                for i, row in enumerate(rows):
                    # Pad row to max_cols and convert None to empty string
                    cells = [str(cell) if cell is not None else "" for cell in row[:max_cols]]
                    while len(cells) < max_cols:
                        cells.append("")
                    table_md.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        table_md.append("| " + " | ".join(["---"] * max_cols) + " |")

                text_content.append("\n".join(table_md))

            workbook.close()
            self.content = "\n\n".join(text_content)
            return True

        except Exception as e:
            _logger.error(f"Error parsing XLSX: {str(e)}")
            self._post_styled_message(f"Error parsing XLSX: {str(e)}", "error")
            return self._parse_default(record, field)

    def _parse_pptx(self, record, field):
        """Parse PowerPoint (.pptx) file and extract text content"""
        if Presentation is None:
            _logger.warning("python-pptx not installed, falling back to default parser")
            return self._parse_default(record, field)

        import io
        pptx_data = field["rawcontent"]

        try:
            pptx_stream = io.BytesIO(pptx_data)
            prs = Presentation(pptx_stream)

            text_content = [f"# {record.name}\n"]

            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"\n## Slide {slide_num}\n")

                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        table_md = []
                        for i, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            table_md.append("| " + " | ".join(cells) + " |")
                            if i == 0:
                                table_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        slide_text.append("\n" + "\n".join(table_md))

                if slide_text:
                    text_content.append("\n".join(slide_text))
                else:
                    text_content.append("*No text content*")

                # Extract notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_content.append(f"\n**Speaker Notes:** {notes}")

            self.content = "\n\n".join(text_content)
            return True

        except Exception as e:
            _logger.error(f"Error parsing PPTX: {str(e)}")
            self._post_styled_message(f"Error parsing PPTX: {str(e)}", "error")
            return self._parse_default(record, field)

    def _parse_csv(self, record, field):
        """Parse CSV file and convert to markdown table"""
        import csv
        import io

        csv_data = field["rawcontent"]

        try:
            # Handle both bytes and string
            if isinstance(csv_data, bytes):
                csv_data = csv_data.decode("utf-8-sig")  # utf-8-sig handles BOM

            csv_stream = io.StringIO(csv_data)

            # Try to detect delimiter
            sample = csv_data[:4096]
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
            except csv.Error:
                dialect = csv.excel  # Default to comma

            csv_stream.seek(0)
            reader = csv.reader(csv_stream, dialect)
            rows = list(reader)

            if not rows:
                self.content = f"# {record.name}\n\n*Empty CSV file*"
                return True

            text_content = [f"# {record.name}\n"]

            # Build markdown table
            table_md = []
            max_cols = max(len(row) for row in rows) if rows else 0

            for i, row in enumerate(rows):
                # Pad row and escape pipe characters
                cells = [cell.replace("|", "\\|") for cell in row]
                while len(cells) < max_cols:
                    cells.append("")
                table_md.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    table_md.append("| " + " | ".join(["---"] * max_cols) + " |")

            text_content.append("\n".join(table_md))
            self.content = "\n\n".join(text_content)
            return True

        except Exception as e:
            _logger.error(f"Error parsing CSV: {str(e)}")
            self._post_styled_message(f"Error parsing CSV: {str(e)}", "error")
            return self._parse_default(record, field)

    def _correct_ocr_text(self, ocr_text):
        """Use LLM to correct OCR errors in Vietnamese/English text"""
        if not ocr_text or len(ocr_text) < 10:
            return ocr_text

        try:
            # Get OCR correction model from collection settings
            ocr_model = None
            provider = None

            for collection in self.collection_ids:
                if collection.ocr_correction_model_id:
                    ocr_model = collection.ocr_correction_model_id
                    provider = ocr_model.provider_id
                    break

            # If no model configured in collection, skip correction
            if not ocr_model or not provider:
                _logger.debug("No OCR correction model configured in collection, skipping correction")
                return ocr_text

            # Prepare prompt for OCR correction
            system_prompt = """Bạn là chuyên gia sửa lỗi văn bản OCR tiếng Việt và tiếng Anh.

NHIỆM VỤ: Sửa lỗi nhận dạng ký tự từ OCR, bao gồm:
- Ký tự bị nhận sai (ví dụ: U thành V, l thành I, 0 thành O)
- Ký tự bị dính liền (ví dụ: "textyaml" → "text/yaml")
- Dấu thanh tiếng Việt bị sai hoặc thiếu
- Lỗi chính tả do OCR

QUY TẮC NGHIÊM NGẶT:
1. CHỈ trả về văn bản đã sửa, KHÔNG giải thích
2. Giữ nguyên cấu trúc, định dạng, xuống dòng
3. KHÔNG thêm hoặc bớt nội dung
4. Nếu thấy từ kỹ thuật như "text/yaml/json", "API", "SSE", "JSON" - hãy sửa về đúng format
5. Với tiếng Việt, chú ý dấu thanh: à á ả ã ạ, è é ẻ ẽ ẹ, etc."""

            user_message = f"""Sửa lỗi OCR cho văn bản sau. Chú ý các từ kỹ thuật IT có thể bị nhận sai ký tự.

VĂN BẢN CẦN SỬA:
{ocr_text}

VĂN BẢN ĐÃ SỬA:"""

            # Call LLM with specific model
            response = provider.chat(
                messages=self.env["mail.message"].browse([]),
                model=ocr_model,
                stream=False,
                prepend_messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message}
                ]
            )

            if response:
                if isinstance(response, dict):
                    corrected = response.get("content", "") or response.get("text", "")
                else:
                    corrected = str(response)

                if corrected and len(corrected) > 5:
                    _logger.info(f"OCR text corrected by {ocr_model.name} (original: {len(ocr_text)} chars, corrected: {len(corrected)} chars)")
                    return corrected.strip()

            return ocr_text

        except Exception as e:
            _logger.warning(f"Failed to correct OCR text with LLM: {e}")
            return ocr_text

    def _parse_image_ocr(self, record, field):
        """Parse image with OCR to extract text"""
        if pytesseract is None or Image is None:
            _logger.warning("pytesseract or PIL not installed, falling back to image reference")
            return self._parse_image(record, field)

        import io
        image_data = field["rawcontent"]

        try:
            image_stream = io.BytesIO(image_data)
            img = Image.open(image_stream)

            # Run OCR with Vietnamese + English languages
            ocr_text = pytesseract.image_to_string(img, lang='vie+eng')

            text_content = [f"# {record.name}\n"]

            # Add image reference
            image_url = f"/web/image/{record.id}"
            text_content.append(f"![{record.name}]({image_url})\n")

            # Add OCR text if found
            if ocr_text.strip():
                # Try to correct OCR text using LLM
                corrected_text = self._correct_ocr_text(ocr_text.strip())
                text_content.append("## Extracted Text (OCR)\n")
                text_content.append(corrected_text)
            else:
                text_content.append("*No text detected in image*")

            self.content = "\n\n".join(text_content)
            return True

        except Exception as e:
            _logger.error(f"Error running OCR: {str(e)}")
            self._post_styled_message(f"Error running OCR: {str(e)}", "warning")
            return self._parse_image(record, field)

    def _parse_default(self, record, field):
        # Default to a generic description for unsupported types
        mimetype = field["mimetype"]
        self.content = f"""
            # {record.name}

            **File Type**: {mimetype}
            **Description**: This file is of type {mimetype} which cannot be directly parsed into text content.
            **Access**: [Open file](/web/content/{record.id})
                            """
        return True
